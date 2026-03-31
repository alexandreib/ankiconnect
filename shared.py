"""Shared utilities for the Anki HSK vocabulary tools."""
import json
import os
import re
import time
import urllib.parse
import urllib.request

# ── Constants ────────────────────────────────────────────────────────────────

DATA_DIR = os.path.join(os.path.dirname(__file__), "data")
CACHE_DIR = os.path.join(os.path.dirname(__file__), "cache")

JSON_FILE = os.path.join(DATA_DIR, "myhsk1_data.json")
POS_CACHE_FILE = os.path.join(CACHE_DIR, "pos_cache.json")
CHAR_INFO_CACHE_FILE = os.path.join(CACHE_DIR, "char_info_cache.json")
NEVER_PROPOSE_FILE = "never_propose_words.txt"
FIXES_FILE = "fixes.json"
PICK_FILE = "pick_words.txt"
NEWWORDS_DIR = "newwords"

POS_KEEP = {"noun", "verb", "adjective", "adverb"}
PRESERVE_TAGS = {"leech", "lesson"}


# ── Utilities ────────────────────────────────────────────────────────────────

def is_cjk(char):
    """Return True if the character is a CJK ideograph."""
    cp = ord(char)
    return (0x4E00 <= cp <= 0x9FFF or 0x3400 <= cp <= 0x4DBF
            or 0x20000 <= cp <= 0x2A6DF or 0xF900 <= cp <= 0xFAFF)


def strip_html(text):
    """Remove HTML tags and &nbsp; entities."""
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'&nbsp;', ' ', text)
    return text.strip()


# ── Google Translate ─────────────────────────────────────────────────────────

def google_translate(word):
    """Full lookup via Google Translate: English, Pinyin, POS tags.
    Raises RuntimeError on failure after 3 retries."""
    encoded = urllib.parse.quote(word)
    url = (
        "https://translate.googleapis.com/translate_a/single"
        f"?client=gtx&sl=zh-CN&tl=en&dt=t&dt=bd&dt=rm&q={encoded}"
    )
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})

    for attempt in range(3):
        try:
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = json.loads(resp.read())
            break
        except Exception as e:
            if attempt < 2:
                time.sleep(2 ** attempt)
            else:
                raise RuntimeError(f"Google Translate failed for '{word}': {e}")

    english = ""
    try:
        english = data[0][0][0]
    except (IndexError, TypeError):
        pass

    pinyin = ""
    try:
        pinyin = data[0][1][3]
    except (IndexError, TypeError):
        pass

    pos_tags = []
    alternatives = []
    if len(data) > 1 and data[1]:
        for entry in data[1]:
            if isinstance(entry, list) and entry:
                pos = entry[0].lower().strip()
                if pos in POS_KEEP and pos not in pos_tags:
                    pos_tags.append(pos)
                # Collect alternative translations
                if len(entry) > 1 and isinstance(entry[1], list):
                    for alt in entry[1]:
                        if isinstance(alt, str):
                            alt = alt.strip()
                            if alt and alt.lower() != english.lower().strip():
                                if alt not in alternatives:
                                    alternatives.append(alt)

    return english, pinyin, pos_tags, alternatives


# ── Cached lookups ───────────────────────────────────────────────────────────

def build_definition(english, alternatives, max_total=5):
    """Combine primary translation with alternatives for a distinctive definition."""
    if not alternatives:
        return english
    parts = [english]
    seen = {english.lower().strip()}
    for alt in alternatives:
        if alt.lower().strip() not in seen:
            seen.add(alt.lower().strip())
            parts.append(alt)
            if len(parts) >= max_total:
                break
    return ", ".join(parts)


def lookup_pos(word, pos_cache):
    """Get POS tags for a word, using cache. Returns (pos_tags, was_cached)."""
    if word in pos_cache:
        return pos_cache[word], True
    try:
        _english, _pinyin, pos_tags, _alts = google_translate(word)
    except RuntimeError as e:
        print(f"  ✗ {e}")
        pos_cache[word] = []
        return [], False
    pos_cache[word] = pos_tags
    return pos_tags, False


def lookup_full(word, pos_cache, char_cache):
    """Get English, Pinyin, POS, alternatives for a word, using char_cache.
    Returns (english, pinyin, pos_tags, alts, was_cached)."""
    if word in char_cache:
        info = char_cache[word]
        return info["english"], info["pinyin"], info["pos"], info.get("alts", []), True
    try:
        english, pinyin, pos_tags, alts = google_translate(word)
    except RuntimeError as e:
        print(f"  ✗ {e}")
        char_cache[word] = {"english": "", "pinyin": "", "pos": [], "alts": []}
        pos_cache[word] = []
        return "", "", [], [], False
    pos_cache[word] = pos_tags
    char_cache[word] = {"english": english, "pinyin": pinyin, "pos": pos_tags, "alts": alts}
    return english, pinyin, pos_tags, alts, False


# ── JSON cache helpers ───────────────────────────────────────────────────────

def load_json_file(filepath):
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {}


def save_json_file(filepath, data):
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def load_pos_cache():
    return load_json_file(POS_CACHE_FILE)


def save_pos_cache(cache):
    save_json_file(POS_CACHE_FILE, cache)


def load_char_info_cache():
    return load_json_file(CHAR_INFO_CACHE_FILE)


def save_char_info_cache(cache):
    save_json_file(CHAR_INFO_CACHE_FILE, cache)


# ── Never-propose list ───────────────────────────────────────────────────────

def load_never_propose():
    """Load characters/words from never_propose_words.txt (one per line, # = comment).
    Supports both plain '的' and '乌 | Wū | black' formats."""
    try:
        with open(NEVER_PROPOSE_FILE, "r", encoding="utf-8") as f:
            entries = set()
            for line in f:
                line = line.strip()
                if line and not line.startswith("#"):
                    entries.add(line.split("|")[0].strip())
            return entries
    except FileNotFoundError:
        return set()


# ── Fixes ────────────────────────────────────────────────────────────────────

def load_fixes():
    """Load definition/pinyin fixes from fixes.json."""
    return load_json_file(FIXES_FILE)


# ── Document extraction ──────────────────────────────────────────────────────

def extract_pptx_text(filepath):
    """Extract all text from a .pptx file."""
    from pptx import Presentation
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:  # type: ignore[attr-defined]
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            texts.append(t)
    return texts


def extract_doc_lines(filepath):
    """Extract text lines from an old .doc (OLE2) file."""
    import olefile
    try:
        ole = olefile.OleFileIO(filepath)
        lines = []
        for sn in ["1Table", "0Table", "WordDocument"]:
            if ole.exists(sn):
                data = ole.openstream(sn).read()
                decoded = data.decode("utf-16-le", errors="ignore")
                for chunk in re.split(r'[\x00-\x08\x0b\x0c\x0e-\x1f]+', decoded):
                    chunk = chunk.strip()
                    if chunk:
                        lines.append(chunk)
        ole.close()
        return lines
    except Exception as e:
        print(f"  Error reading {filepath}: {e}")
        return []


# ── Deck data ────────────────────────────────────────────────────────────────

def load_deck_records():
    """Load all records from JSON_FILE."""
    with open(JSON_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def save_deck_records(records):
    """Save records to JSON_FILE."""
    save_json_file(JSON_FILE, records)


def load_existing_words():
    """Load existing Chinese words and chars from the deck JSON.
    Returns (words_set, chars_set)."""
    data = load_deck_records()
    words = set()
    chars = set()
    for note in data:
        w = strip_html(note["fields"].get("中文", ""))
        if w:
            words.add(w)
            for c in w:
                if is_cjk(c):
                    chars.add(c)
    return words, chars
